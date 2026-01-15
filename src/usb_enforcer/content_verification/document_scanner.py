"""
Document format scanner.

Extracts and scans text from various document formats including
PDF, DOCX, XLSX, PPTX, and ODF formats.
"""

import logging
from pathlib import Path
from typing import Optional


logger = logging.getLogger(__name__)


class DocumentScanner:
    """
    Scanner for various document formats.
    
    Extracts text content from office documents and PDFs for
    sensitive data scanning.
    """
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.xlsx': 'xlsx',
        '.pptx': 'pptx',
        '.odt': 'odt',
        '.ods': 'ods',
        '.odp': 'odp',
        '.doc': 'doc',  # Old Word format
        '.xls': 'xls',  # Old Excel format
        '.ppt': 'ppt',  # Old PowerPoint format
        '.msg': 'msg',  # Outlook message
        '.rtf': 'rtf',  # Rich Text Format
    }
    
    def __init__(self, content_scanner):
        """
        Initialize document scanner.
        
        Args:
            content_scanner: ContentScanner instance for scanning extracted text
        """
        self.content_scanner = content_scanner
        
        logger.info("Document scanner initialized")
    
    def is_document(self, filepath: Path) -> bool:
        """
        Check if file is a supported document format.
        
        Args:
            filepath: Path to check
            
        Returns:
            True if file is a supported document
        """
        return filepath.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def scan_document(self, filepath: Path):
        """
        Scan document for sensitive content.
        
        Args:
            filepath: Path to document file
            
        Returns:
            ScanResult from content scanner
        """
        from .scanner import ScanResult, ScanAction
        
        logger.debug(f"Scanning document: {filepath.name}")
        
        doc_type = self.SUPPORTED_EXTENSIONS.get(filepath.suffix.lower())
        
        if not doc_type:
            return ScanResult(
                blocked=False,
                action=ScanAction.ALLOW,
                reason="Unsupported document type"
            )
        
        try:
            # Extract text based on document type
            if doc_type == 'pdf':
                text = self._extract_pdf_text(filepath)
            elif doc_type == 'docx':
                text = self._extract_docx_text(filepath)
            elif doc_type == 'xlsx':
                text = self._extract_xlsx_text(filepath)
            elif doc_type == 'pptx':
                text = self._extract_pptx_text(filepath)
            elif doc_type in ('odt', 'ods', 'odp'):
                text = self._extract_odf_text(filepath, doc_type)
            elif doc_type == 'doc':
                text = self._extract_doc_text(filepath)
            elif doc_type == 'xls':
                text = self._extract_xls_text(filepath)
            elif doc_type == 'ppt':
                text = self._extract_ppt_text(filepath)
            elif doc_type == 'msg':
                text = self._extract_msg_text(filepath)
            elif doc_type == 'rtf':
                text = self._extract_rtf_text(filepath)
            else:
                return ScanResult(
                    blocked=False,
                    action=ScanAction.ALLOW,
                    reason=f"No handler for {doc_type}"
                )
            
            if not text:
                return ScanResult(
                    blocked=False,
                    action=ScanAction.ALLOW,
                    reason="No text content extracted"
                )
            
            # Scan extracted text
            result = self.content_scanner.scan_content(text.encode('utf-8'), filepath.name)
            
            if result.blocked:
                result.location = f"{filepath.name} ({doc_type})"
            
            return result
            
        except Exception as e:
            logger.error(f"Error scanning document {filepath.name}: {e}", exc_info=True)
            
            if self.content_scanner.config.get('block_on_error', True):
                return ScanResult(
                    blocked=True,
                    action=ScanAction.BLOCK,
                    reason=f"Document scan error: {str(e)}"
                )
            else:
                return ScanResult(
                    blocked=False,
                    action=ScanAction.ALLOW,
                    reason=f"Document scan error (allowed): {str(e)}"
                )
    
    def _extract_pdf_text(self, filepath: Path) -> str:
        """Extract text from PDF"""
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber not installed, skipping PDF")
            return ""
        
        try:
            text_parts = []
            
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    # Extract text
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                    
                    # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        for row in table:
                            if row:
                                row_text = ' '.join(str(cell) for cell in row if cell)
                                text_parts.append(row_text)
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _extract_docx_text(self, filepath: Path) -> str:
        """Extract text from DOCX"""
        try:
            import docx
        except ImportError:
            logger.warning("python-docx not installed, skipping DOCX")
            return ""
        
        try:
            doc = docx.Document(filepath)
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_parts.append(cell.text)
            
            # Extract headers/footers
            for section in doc.sections:
                if section.header:
                    for paragraph in section.header.paragraphs:
                        if paragraph.text:
                            text_parts.append(paragraph.text)
                
                if section.footer:
                    for paragraph in section.footer.paragraphs:
                        if paragraph.text:
                            text_parts.append(paragraph.text)
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _extract_xlsx_text(self, filepath: Path) -> str:
        """Extract text from XLSX"""
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl not installed, skipping XLSX")
            return ""
        
        try:
            workbook = openpyxl.load_workbook(filepath, data_only=True)
            text_parts = []
            
            for sheet in workbook.worksheets:
                # Add sheet name
                text_parts.append(f"Sheet: {sheet.title}")
                
                # Extract cell values
                for row in sheet.iter_rows(values_only=True):
                    for cell_value in row:
                        if cell_value is not None:
                            text_parts.append(str(cell_value))
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {e}")
            return ""
    
    def _extract_pptx_text(self, filepath: Path) -> str:
        """Extract text from PPTX"""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed, skipping PPTX")
            return ""
        
        try:
            prs = Presentation(filepath)
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text:
                            text_parts.append(shape.text)
                    
                    # Extract table text
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_parts.append(cell.text)
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_odf_text(self, filepath: Path, doc_type: str) -> str:
        """Extract text from ODF formats (ODT, ODS, ODP)"""
        try:
            from odf import text, table, teletype
            from odf.opendocument import load
        except ImportError:
            logger.warning("odfpy not installed, skipping ODF")
            return ""
        
        try:
            textdoc = load(filepath)
            text_parts = []
            
            # Extract text elements
            for element in textdoc.getElementsByType(text.P):
                text_content = teletype.extractText(element)
                if text_content:
                    text_parts.append(text_content)
            
            # Extract table content for ODS
            if doc_type == 'ods':
                for tab in textdoc.spreadsheet.getElementsByType(table.Table):
                    for row in tab.getElementsByType(table.TableRow):
                        for cell in row.getElementsByType(table.TableCell):
                            cell_text = teletype.extractText(cell)
                            if cell_text:
                                text_parts.append(cell_text)
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting ODF text: {e}")
            return ""
    
    def _extract_doc_text(self, filepath: Path) -> str:
        """Extract text from old Word .doc format (OLE2)"""
        try:
            import olefile
        except ImportError:
            logger.warning("olefile not installed, skipping .doc")
            return ""
        
        try:
            # Basic text extraction from OLE2 Word document
            # This is limited - .doc format is complex
            ole = olefile.OleFileIO(filepath)
            
            # Try to extract WordDocument stream
            if ole.exists('WordDocument'):
                # This is a simplified approach - full parsing is complex
                # For production, might want antiword or similar
                logger.debug(f"Detected Word .doc file: {filepath.name}")
                # For now, just mark as detected but can't extract
                ole.close()
                return ""  # Would need antiword or more complex parsing
            
            ole.close()
            return ""
            
        except Exception as e:
            logger.error(f"Error extracting .doc text: {e}")
            return ""
    
    def _extract_xls_text(self, filepath: Path) -> str:
        """Extract text from old Excel .xls format"""
        try:
            import xlrd
        except ImportError:
            logger.warning("xlrd not installed, skipping .xls")
            return ""
        
        try:
            workbook = xlrd.open_workbook(filepath)
            text_parts = []
            
            for sheet in workbook.sheets():
                for row_idx in range(sheet.nrows):
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            text_parts.append(str(cell.value))
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting .xls text: {e}")
            return ""
    
    def _extract_ppt_text(self, filepath: Path) -> str:
        """Extract text from old PowerPoint .ppt format (OLE2)"""
        try:
            import olefile
        except ImportError:
            logger.warning("olefile not installed, skipping .ppt")
            return ""
        
        try:
            # Basic detection of OLE2 PowerPoint
            # Full text extraction would require complex parsing
            ole = olefile.OleFileIO(filepath)
            
            if ole.exists('PowerPoint Document'):
                logger.debug(f"Detected PowerPoint .ppt file: {filepath.name}")
                # Complex format - would need specialized parser
                ole.close()
                return ""
            
            ole.close()
            return ""
            
        except Exception as e:
            logger.error(f"Error extracting .ppt text: {e}")
            return ""
    
    def _extract_msg_text(self, filepath: Path) -> str:
        """Extract text from Outlook .msg format"""
        try:
            import extract_msg
        except ImportError:
            logger.warning("extract-msg not installed, skipping .msg")
            return ""
        
        try:
            msg = extract_msg.Message(filepath)
            text_parts = []
            
            # Extract subject, sender, body
            if msg.subject:
                text_parts.append(f"Subject: {msg.subject}")
            if msg.sender:
                text_parts.append(f"From: {msg.sender}")
            if msg.body:
                text_parts.append(msg.body)
            
            # Extract attachment names (but not content - handled by archive scanner)
            if hasattr(msg, 'attachments'):
                for attachment in msg.attachments:
                    if hasattr(attachment, 'longFilename'):
                        text_parts.append(f"Attachment: {attachment.longFilename}")
            
            msg.close()
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting .msg text: {e}")
            return ""
    
    def _extract_rtf_text(self, filepath: Path) -> str:
        """Extract text from RTF format"""
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            logger.warning("striprtf not installed, skipping RTF")
            return ""
        
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            return text
            
        except Exception as e:
            logger.error(f"Error extracting RTF text: {e}")
            return ""

